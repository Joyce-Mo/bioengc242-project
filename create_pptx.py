#!/usr/bin/env python
"""Generate project update PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_MED = RGBColor(0x16, 0x21, 0x3E)
ACCENT = RGBColor(0x53, 0x3C, 0xA6)
ACCENT2 = RGBColor(0x6C, 0x63, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
TEAL = RGBColor(0x00, 0xD2, 0xD3)
CORAL = RGBColor(0xFF, 0x6B, 0x6B)
GREEN = RGBColor(0x1D, 0xD1, 0xA1)
YELLOW = RGBColor(0xFE, 0xCA, 0x57)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_slide(slide, left, top, width, height, items, font_size=16,
                     color=WHITE, bullet_color=TEAL):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)
        # Bullet character
        run_bullet = p.add_run()
        run_bullet.text = "\u25B8 "
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = "Calibri"
        # Text
        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = "Calibri"
    return tf


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, items, accent_color=ACCENT2,
             font_size=14):
    add_rounded_rect(slide, left, top, width, height, RGBColor(0x1F, 0x2B, 0x4D), accent_color)
    # Title bar
    add_rounded_rect(slide, left, top, width, 0.45, accent_color)
    add_textbox(slide, left + 0.15, top + 0.05, width - 0.3, 0.4, title,
                font_size=font_size + 2, color=WHITE, bold=True)
    add_bullet_slide(slide, left + 0.1, top + 0.5, width - 0.2, height - 0.6,
                     items, font_size=font_size, color=LIGHT_GRAY)


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

add_textbox(slide, 1.5, 1.5, 10, 1.2,
            "Mapping Multi-Conformer Protein Ensembles\nwith Diffusion Models",
            font_size=38, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 3.2, 10, 0.6,
            "BioEng C242 Project Update",
            font_size=24, color=ACCENT2, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 4.2, 10, 0.5,
            "Joyce Mo  |  March 2026",
            font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Decorative line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(4), Inches(4.0), Inches(5.333), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT2
shape.line.fill.background()

# ============================================================
# SLIDE 2: Project Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_textbox(slide, 0.8, 0.4, 11, 0.7, "Project Overview", font_size=32, color=WHITE, bold=True)
# Underline
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0.8), Inches(1.0), Inches(2.5), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT2
shape.line.fill.background()

add_textbox(slide, 0.8, 1.3, 11.5, 0.8,
            "Goal: Develop generative models that capture the conformational diversity of proteins,\n"
            "moving beyond single-structure prediction to multi-conformer ensemble generation.",
            font_size=18, color=LIGHT_GRAY)

add_card(slide, 0.8, 2.5, 3.7, 3.8, "Motivation", [
    "Proteins are dynamic, not static",
    "Single structures miss functional states",
    "Ensembles capture flexibility, allostery,\n  and binding diversity",
    "Current methods (AlphaFold) predict\n  one conformation",
], accent_color=ACCENT2, font_size=14)

add_card(slide, 4.8, 2.5, 3.7, 3.8, "Approach", [
    "Two novel architectures:",
    "  1. VAE + Discrete Diffusion (DMRA)",
    "  2. Flow-Matching on all-atom coords",
    "Train on synthetic ensembles from\n  MC-SCE and Rosetta Backrub",
    "CATH20 dataset for structural diversity",
], accent_color=TEAL, font_size=14)

add_card(slide, 8.8, 2.5, 3.7, 3.8, "Key Innovation", [
    "Representation alignment bridges\n  VAE latent space and diffusion",
    "Ensemble consistency loss ensures\n  coherent multi-conformer output",
    "Flow-matching provides deterministic\n  ODE-based sampling",
], accent_color=CORAL, font_size=14)

# ============================================================
# SLIDE 3: Data Pipeline
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_textbox(slide, 0.8, 0.4, 11, 0.7, "Data Pipeline", font_size=32, color=WHITE, bold=True)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0.8), Inches(1.0), Inches(2.0), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT2
shape.line.fill.background()

# Pipeline boxes
stages = [
    ("Raw PDB Files\n(CATH20)", ACCENT2, 0.8),
    ("MC-SCE /\nBackrub", TEAL, 3.3),
    ("Ensemble\nPDB Files", GREEN, 5.8),
    ("Graph\nBuilder", YELLOW, 8.3),
    ("PyG Data\n(.pt files)", CORAL, 10.8),
]

for label, color, left in stages:
    add_rounded_rect(slide, left, 1.5, 2.0, 1.2, RGBColor(0x1F, 0x2B, 0x4D), color)
    add_textbox(slide, left + 0.1, 1.6, 1.8, 1.0, label,
                font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Arrows between stages
for left in [2.8, 5.3, 7.8, 10.3]:
    add_textbox(slide, left, 1.7, 0.5, 0.5, "\u279C", font_size=28, color=ACCENT2,
                alignment=PP_ALIGN.CENTER)

# Feature details
add_card(slide, 0.8, 3.2, 3.7, 3.5, "Node Features (12-dim)", [
    "B-factor (normalized)",
    "SASA (solvent accessibility)",
    "Backbone dihedrals: sin/cos(\u03C6, \u03C8)",
    "AA properties: hydrophobicity,\n  charge, polarity, volume",
    "Local density (10\u00C5 radius)",
    "Sequence-structure exposure",
], accent_color=ACCENT2, font_size=13)

add_card(slide, 4.8, 3.2, 3.7, 3.5, "Edge Features (93-dim)", [
    "Gaussian RBF distance (15-dim)",
    "Relative backbone positions (12-dim)",
    "Sequence distance encoding (66-dim)",
    "K-NN graph (K=30 on C\u03B1)",
], accent_color=TEAL, font_size=13)

add_card(slide, 8.8, 3.2, 3.7, 3.5, "Ensemble Generation", [
    "MC-SCE: Side-chain rotamer sampling\n  at T=300K, 5 conformers/protein",
    "Backrub: 10K MC steps at kT=0.6\n  75% backbone / 25% sidechain moves",
    "Running on Wynton (SGE)\n  and Savio (SLURM) clusters",
], accent_color=GREEN, font_size=13)

# ============================================================
# SLIDE 4: Architecture 1 - VAE + Discrete Diffusion
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_textbox(slide, 0.8, 0.4, 11, 0.7, "Architecture 1: VAE + Discrete Diffusion (DMRA)",
            font_size=30, color=WHITE, bold=True)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0.8), Inches(1.0), Inches(5.0), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT2
shape.line.fill.background()

# VAE card
add_card(slide, 0.8, 1.3, 3.7, 2.8, "VAE (Structural Encoder)", [
    "4-layer GNN encoder \u2192 latent \u03BC, \u03C3",
    "Latent dim: 128, Hidden: 256",
    "Per-node embeddings for alignment",
    "4-layer GNN decoder reconstructs\n  12-dim node features",
    "\u03B2-annealing: 0 \u2192 1 over 20 epochs",
], accent_color=ACCENT2, font_size=13)

# Diffusion card
add_card(slide, 4.8, 1.3, 3.7, 2.8, "Discrete Diffusion", [
    "Categorical diffusion on 20 AA types",
    "Forward: Q_t = \u03B1_t\u00B7I + (1-\u03B1_t)/20\u00B711\u1D40",
    "Cosine noise schedule, T=500 steps",
    "Per-residue independent transitions",
    "Ancestral sampling for denoising",
], accent_color=TEAL, font_size=13)

# Denoiser card
add_card(slide, 8.8, 1.3, 3.7, 2.8, "DMRA Denoiser", [
    "6 layers of edge-conditioned MP",
    "Multi-head attention (8 heads)",
    "Shared center attention (virtual\n  global node, O(n) complexity)",
    "VAE alignment injection every\n  2 layers via gated fusion",
    "Output: logits over 20 AA types",
], accent_color=CORAL, font_size=13)

# Alignment section
add_card(slide, 0.8, 4.4, 5.7, 2.5, "Representation Alignment", [
    "VAE embeddings \u2192 shared alignment space \u2190 Denoiser embeddings",
    "Cosine similarity loss bridges the two representation spaces",
    "Gated injection: h = h_denoiser + \u03C3(gate) \u00B7 norm(transform(concat(h_denoiser, h_vae)))",
    "Enables structural conditioning of the diffusion process",
], accent_color=YELLOW, font_size=14)

# Training info
add_card(slide, 6.8, 4.4, 5.7, 2.5, "Training Protocol", [
    "Phase 1: VAE pre-training (50 epochs, cosine LR)",
    "Phase 2: Joint training (100 epochs)",
    "  L = L_diffusion + \u03BB_vae\u00B7L_vae + \u03BB_align\u00B7L_align",
    "  VAE frozen for first 5 epochs, then unfrozen",
    "  \u03BB_vae=1.0, \u03BB_align=0.5, LR=3e-5",
], accent_color=GREEN, font_size=14)


# ============================================================
# SLIDE 5: Architecture 2 - Flow Matching
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_textbox(slide, 0.8, 0.4, 11, 0.7, "Architecture 2: Flow-Matching Diffusion",
            font_size=30, color=WHITE, bold=True)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0.8), Inches(1.0), Inches(4.5), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = TEAL
shape.line.fill.background()

add_card(slide, 0.8, 1.3, 5.7, 2.8, "Flow-Matching Framework", [
    "Linear interpolation path: x_t = (1-t)\u00B7x_0 + t\u00B7\u03B5,  t \u2208 [0,1]",
    "Learn velocity field: v = \u03B5 - x_0",
    "Deterministic ODE integration (Euler method, 200 steps)",
    "Operates on all-atom coordinates (atom37 representation)",
    "  \u2022 Positions 0-3: N, CA, C, O backbone atoms",
    "  \u2022 Position 4: C\u03B2 (ideal geometry, masked for GLY)",
    "Loss: ||v_pred - (\u03B5 - x_0)||\u00B2 with per-atom masking",
], accent_color=TEAL, font_size=14)

add_card(slide, 6.8, 1.3, 5.7, 2.8, "CoordinateVelocityNet", [
    "Input: noisy atom37 coords + atom mask",
    "Time conditioning: sinusoidal embeddings",
    "Self-conditioning from previous prediction",
    "12 layers of edge-conditioned message passing",
    "Shared center attention for long-range info",
    "K-NN graph rebuilt from noisy C\u03B1 positions",
    "Output: per-atom velocity field (37\u00D73 per residue)",
], accent_color=ACCENT2, font_size=14)

# Comparison
add_card(slide, 0.8, 4.5, 11.7, 2.3, "Architecture Comparison", [
    "VAE + Discrete Diffusion:  Operates on amino acid types (categorical)  |  Stochastic sampling  |  "
    "Conditioned via VAE latent space  |  Suited for sequence design",
    "Flow-Matching:  Operates on 3D coordinates (continuous)  |  Deterministic ODE  |  "
    "Direct structural generation  |  Suited for conformer generation",
    "Both use edge-conditioned message passing + shared center attention as core GNN building blocks",
], accent_color=CORAL, font_size=14)

# ============================================================
# SLIDE 6: GNN Building Blocks
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_textbox(slide, 0.8, 0.4, 11, 0.7, "Core GNN Building Blocks",
            font_size=30, color=WHITE, bold=True)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0.8), Inches(1.0), Inches(3.5), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT2
shape.line.fill.background()

add_card(slide, 0.8, 1.3, 3.7, 3.0, "Edge-Conditioned\nMessage Passing", [
    "Multi-head attention (8 heads)",
    "Score = (q\u00B7k)/\u221Ad + edge_bias",
    "Edge features \u2192 per-head bias via MLP",
    "Feed-forward net + residual connection",
    "LayerNorm + dropout (0.1)",
    "Used in: VAE, Denoiser, Flow-Matching",
], accent_color=ACCENT2, font_size=14)

add_card(slide, 4.8, 1.3, 3.7, 3.0, "Shared Center Attention", [
    "Virtual global node for O(n) attention",
    "Phase 1: Residues \u2192 center (attention)",
    "Phase 2: Center \u2192 residues (gating)",
    "Replaces O(n\u00B2) full self-attention",
    "Per-graph via to_dense_batch",
    "Used in: Denoiser, Flow-Matching",
], accent_color=TEAL, font_size=14)

add_card(slide, 8.8, 1.3, 3.7, 3.0, "K-NN Graph Construction", [
    "K=30 neighbors on C\u03B1 atoms",
    "93-dim edge features encode:",
    "  \u2022 Euclidean distance (RBF)",
    "  \u2022 Relative backbone geometry",
    "  \u2022 Sequence separation",
    "Rebuilt per-sample (or per-step\n  in flow-matching)",
], accent_color=GREEN, font_size=14)

add_card(slide, 0.8, 4.6, 5.7, 2.3, "Loss Functions", [
    "VAE Loss: MSE reconstruction + \u03B2\u00B7KL divergence (\u03B2-annealing)",
    "Diffusion Loss: Cross-entropy on predicted clean AA type logits",
    "Alignment Loss: Cosine similarity between VAE and denoiser embeddings",
    "Ensemble Consistency: KL divergence between conformer latent distributions",
    "Flow Loss: MSE on predicted velocity field with atom masking",
], accent_color=YELLOW, font_size=14)

add_card(slide, 6.8, 4.6, 5.7, 2.3, "Model Configuration", [
    "Hidden dim: 256  |  Latent dim: 128",
    "Num layers: 6 (denoiser/VAE), 12 (flow-matching)",
    "Attention heads: 8  |  Dropout: 0.1",
    "Diffusion timesteps: 500 (cosine schedule)",
    "Flow-matching: 200 ODE steps, \u03C3_data=10.0",
], accent_color=ACCENT2, font_size=14)

# ============================================================
# SLIDE 7: Current Progress
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_textbox(slide, 0.8, 0.4, 11, 0.7, "Current Progress", font_size=32, color=WHITE, bold=True)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0.8), Inches(1.0), Inches(2.5), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = GREEN
shape.line.fill.background()

# Completed
add_card(slide, 0.8, 1.3, 3.7, 5.2, "Completed", [
    "Full model implementation:",
    "  \u2022 VAE encoder/decoder",
    "  \u2022 Discrete diffusion",
    "  \u2022 DMRA denoiser",
    "  \u2022 Flow-matching model",
    "  \u2022 All loss functions",
    "",
    "Data pipeline:",
    "  \u2022 PDB parser + featurizer",
    "  \u2022 Graph builder + datasets",
    "  \u2022 Preprocessing scripts",
    "",
    "Training infrastructure:",
    "  \u2022 Phase 1 & 2 training scripts",
    "  \u2022 Flow-matching training",
    "  \u2022 HPC job submission (SGE/SLURM)",
], accent_color=GREEN, font_size=13)

# In progress
add_card(slide, 4.8, 1.3, 3.7, 5.2, "In Progress", [
    "Ensemble generation:",
    "  \u2022 MC-SCE running on Wynton",
    "  \u2022 Backrub (PyRosetta) on Wynton",
    "  \u2022 30-protein ai-cath subset",
    "",
    "HPC debugging:",
    "  \u2022 Savio CPU allocation issues",
    "  \u2022 Wynton conda env / module",
    "    loading in SGE jobs",
    "",
    "Dataset preparation:",
    "  \u2022 Generating PDB file lists",
    "  \u2022 Path mapping across clusters",
], accent_color=YELLOW, font_size=13)

# Not started
add_card(slide, 8.8, 1.3, 3.7, 5.2, "Not Yet Started", [
    "Model training:",
    "  \u2022 Phase 1: VAE pre-training",
    "  \u2022 Phase 2: Joint training",
    "  \u2022 Flow-matching training",
    "",
    "Evaluation:",
    "  \u2022 Ensemble quality metrics",
    "  \u2022 RMSD / TM-score analysis",
    "  \u2022 Comparison to ground truth",
    "",
    "Analysis & visualization:",
    "  \u2022 Latent space exploration",
    "  \u2022 Generated ensemble diversity",
], accent_color=CORAL, font_size=13)

# ============================================================
# SLIDE 8: Next Steps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_textbox(slide, 0.8, 0.4, 11, 0.7, "Next Steps", font_size=32, color=WHITE, bold=True)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0.8), Inches(1.0), Inches(1.8), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = CORAL
shape.line.fill.background()

add_card(slide, 0.8, 1.3, 3.7, 2.5, "Immediate (This Week)", [
    "Complete MC-SCE + Backrub ensemble\n  generation on Wynton",
    "Resolve Savio allocation issue\n  (contact BRC support)",
    "Preprocess all ensembles to .pt graphs",
    "Begin Phase 1 VAE pre-training",
], accent_color=CORAL, font_size=14)

add_card(slide, 4.8, 1.3, 3.7, 2.5, "Short-Term (1-2 Weeks)", [
    "Complete VAE pre-training (50 epochs)",
    "Run Phase 2 joint training (100 epochs)",
    "Train flow-matching model (100 epochs)",
    "Monitor training with W&B logging",
], accent_color=YELLOW, font_size=14)

add_card(slide, 8.8, 1.3, 3.7, 2.5, "Medium-Term (3-4 Weeks)", [
    "Evaluate generated ensembles:",
    "  \u2022 Structural validity (bond lengths)",
    "  \u2022 Diversity (pairwise RMSD)",
    "  \u2022 Coverage (TM-score to targets)",
    "Compare VAE+Diffusion vs Flow-Matching",
], accent_color=GREEN, font_size=14)

add_card(slide, 0.8, 4.2, 5.7, 2.8, "Key Risks & Mitigations", [
    "HPC allocation: Savio ic_chem242 account blocked \u2192 fallback to Wynton",
    "Training compute: GPU availability on Savio \u2192 request savio3_gpu allocation",
    "Ensemble quality: Synthetic data may not capture real dynamics \u2192 validate against\n"
    "  experimental NMR ensembles if available",
    "Model convergence: Complex joint training \u2192 staged approach with frozen VAE warmup",
], accent_color=ACCENT2, font_size=14)

add_card(slide, 6.8, 4.2, 5.7, 2.8, "Stretch Goals", [
    "Scale to full CATH20 dataset (13K+ proteins)",
    "Incorporate experimental ensemble data (NMR, MD)",
    "Conditional generation: ensemble conditioned on\n  target function or binding partner",
    "Benchmark against existing methods (EnsembleFold,\n  AlphaFold ensembles, MD simulations)",
], accent_color=TEAL, font_size=14)


# Save
out_path = "/Users/joycemo/Documents/GitHub/bioengc242-project/project_update.pptx"
prs.save(out_path)
print(f"Saved to {out_path}")
